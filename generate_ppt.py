from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Helper function to add slides
def add_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = title_text
    tf = content.text_frame
    
    for i, pt in enumerate(bullet_points):
        if i == 0:
            tf.text = pt
        else:
            p = tf.add_paragraph()
            p.text = pt
            p.level = 0

# Slide 1: Title
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "ShopEase AI Customer Support Assistant"
subtitle.text = "Scaling E-Commerce Support with RAG\nPresented by: The Engineering Team"

# Slide 2: The Problem
add_slide("The Challenge in E-Commerce Support", [
    "High Volume: Support teams are overwhelmed by repetitive questions.",
    "High Cost: Scaling human support linearly with business growth is expensive.",
    "Customer Frustration: Customers want instant answers 24/7.",
    "The AI Gap: Generic AI chatbots hallucinate policies and frustrate users."
])

# Slide 3: Our Solution
add_slide("Intelligent RAG with Human-in-the-Loop", [
    "Grounding in Truth (RAG): AI answers using our ingested company documents.",
    "Seamless Escalation (HITL): If the AI doesn't know the answer, it escalates.",
    "Safety Checks: Detects high-risk keywords (e.g., 'fraud', 'lawsuit').",
    "Agent Dashboard: Human agents get a clean dashboard with the escalated ticket."
])

# Slide 4: Key Features
add_slide("What Makes ShopEase AI Powerful?", [
    "Dynamic Knowledge Base: Easily upload PDFs, text files, and Markdown.",
    "Smart Intent Classification: Categorises queries before deciding how to handle them.",
    "Fast & Local Search: Uses ChromaDB for offline vector search.",
    "Cost-Effective LLM: Powered by Groq for ultra-fast, low-cost AI inference."
])

# Slide 5: System Architecture
add_slide("A Robust, Full-Stack Pipeline", [
    "Frontend: User chats via a clean web interface.",
    "FastAPI Backend: Receives and processes the query.",
    "LangGraph State Machine: Orchestrates the AI flow.",
    "ChromaDB: Retrieves relevant policy documents.",
    "Groq (LLM): Generates the final grounded answer.",
    "HITL Queue: Stores unresolved queries as JSON for human review."
])

# Slide 6: How It Works - The Pipeline
add_slide("Under the Hood (LangGraph State Machine)", [
    "1. Input: Validate and classify the user's query.",
    "2. Retrieval: Search the vector database for relevant policies.",
    "3. Router: Decide whether to generate an answer or escalate.",
    "4. Generate: Send the prompt and retrieved documents to the LLM.",
    "5. Escalate (HITL): If confidence is low, create a support ticket.",
    "6. Output: Return the answer with source citations."
])

# Slide 7: Intelligent Escalation
add_slide("When does the AI hand off to a human?", [
    "Hard Keywords: Detects legal or fraud-related terms.",
    "Intent: Automatically escalates complaints.",
    "Low Confidence: If the vector database returns no relevant documents.",
    "LLM Admission: If the LLM generates a response saying 'I don't have enough information'."
])

# Slide 8: The Technology Stack
add_slide("Modern, Scalable, and Fast", [
    "Backend: Python, FastAPI, Uvicorn",
    "AI Orchestration: LangGraph, LangChain",
    "LLM Provider: Groq (llama-3.1-8b-instant)",
    "Vector Database: ChromaDB",
    "Document Processing: PyMuPDF, Pure-NumPy TF-IDF Embedder",
    "Frontend: Vanilla HTML/JS/CSS (Lightweight and fast)"
])

# Slide 9: Next Steps
add_slide("Demo & Next Steps", [
    "Demo 1: Standard policy question (handled by AI).",
    "Demo 2: Out-of-scope question (escalates to Agent Dashboard).",
    "Next Steps:",
    " - Integrate with real customer databases (Shopify/Magento).",
    " - Swap TF-IDF for advanced Neural Embeddings.",
    " - Deploy to production environment."
])

prs.save("ShopEase_Presentation.pptx")
print("Presentation generated successfully!")
